# Значения
# Рамки
# Пробелы по разрядам
# ПНР
# КБО
# Заливка
# Норма
# Ширина / высота столбцов
# Серая зона
# Стрелки

import pptx
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx import Presentation
from datetime import datetime

date = datetime.now().date()
PasteDate = f'{date.day}.{date.month}.{date.year}'
Name = f'ДЗ ЧГБ на {PasteDate}.pptx'

prs = Presentation('###.pptx')
slide = prs.slides[0]

text_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.25), Inches(11.5), Inches(0.5))  # Место
text_box.text = f"Предварительный анализ ДЗ филиалов на {PasteDate}"  # Текст
text_box.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
text_box.text_frame.paragraphs[0].runs[0].font.size = Pt(20)  # Размер шрифта
text_box.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирность
text_box.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет

text_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.7), Inches(3), Inches(0.4))  # Место
text_box.text = f"ВДГО и ВДС ИЖС"  # Текст
text_box.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
text_box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)  # Размер шрифта
text_box.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирность
text_box.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(129, 14, 14)  # Цвет

text_box = slide.shapes.add_textbox(Inches(2.7), Inches(0.76), Inches(2), Inches(0.4))  # Место
text_box.text = f"Данные в таблице из BI"  # Текст
text_box.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
text_box.text_frame.paragraphs[0].runs[0].font.size = Pt(12)  # Размер шрифта
text_box.text_frame.paragraphs[0].runs[0].font.bold = False  # Жирность
text_box.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(129, 14, 14)  # Цвет

text_box = slide.shapes.add_textbox(Inches(0.2), Inches(6.9), Inches(12), Inches(0.6))  # Место
text_box.text = f"Сведения ежемесячно представляются на ЧГБ по закрытому периоду, который был завершен\nв отчетном месяце"  # Текст
for i in range(2):
    text_box.text_frame.paragraphs[i].runs[0].font.name = 'Sitka Small'  # Шрифт
    text_box.text_frame.paragraphs[i].runs[0].font.size = Pt(16)  # Размер шрифта
    text_box.text_frame.paragraphs[i].runs[0].font.bold = True  # Жирность
    text_box.text_frame.paragraphs[i].runs[0].font.color.rgb = pptx.dml.color.RGBColor(129, 14, 14)  # Цвет

rows = 11
cols = 7
table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.1), Inches(12.88), Inches(5.7)).table

for i in range(10):
    for y in range(7):
        cell = table.cell(i, y)
        cell.text = f'{i}:{y}'
        cell.fill.background()  # Заливка
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
        cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Размер шрифта
        cell.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирный
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

cell = table.cell(1, 3)
cell.text = 'КЗ МОГС'
cell.fill.background()  # Заливка
cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Размер шрифта
cell.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирный
cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

cell = table.cell(1, 4)
cell.text = 'ДЗ МОГ'
cell.fill.background()  # Заливка
cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(255, 0, 0)  # Цвет
cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Размер шрифта
cell.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирный
cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

for i in range(5):
    cell = table.cell(2, i + 2)
    cell.text = f'{i + 1}'
    cell.fill.background()  # Заливка
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(112, 48, 160)  # Цвет
    cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)  # Размер шрифта
    cell.text_frame.paragraphs[0].runs[0].font.bold = False  # Жирный
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

txt = ['тыс. руб.', 'тыс. руб.', 'тыс. руб.', 'тыс. руб.', '%']
for i in range(5):
    cell = table.cell(3, i + 2)
    cell.text = txt[i]
    cell.fill.background()  # Заливка
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(112, 48, 160)  # Цвет
    cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)  # Размер шрифта
    cell.text_frame.paragraphs[0].runs[0].font.bold = False  # Жирный
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

for i in range(6):
    cell = table.cell(i + 4, 0)
    cell.text = f'{i + 1}'
    cell.fill.background()  # Заливка
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
    cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)  # Размер шрифта
    cell.text_frame.paragraphs[0].runs[0].font.bold = False  # Жирный
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

Branch = ['Восток', 'Запад', 'Север', 'СЗ', 'Юг', 'ЮВ']
Sort = []
for i in range(6):
    All = 999 * (i + 50)
    KZ = 666
    DZ = All - KZ
    Max = 333
    Itog = round(DZ / Max, 0)
    txt = [f'{Branch[i]}', f'{All}', f'{KZ}', f'{DZ}', f'{Max}', f'{Itog}']
    Sort.append(txt)

for i in range(len(Sort)):
    # Исходно считаем наименьшим первый элемент
    lowest_value_index = i
    # Этот цикл перебирает несортированные элементы
    for j in range(i + 1, len(Sort)):
        if Sort[j][5] > Sort[lowest_value_index][5]:
            lowest_value_index = j
    # Самый маленький элемент меняем с первым в списке
    Sort[i], Sort[lowest_value_index] = Sort[lowest_value_index], Sort[i]

for i in range(6):
    for y in range(6):
        cell = table.cell(i + 4, y + 1)
        cell.text = f'{Sort[i][y]}'
        cell.fill.background()  # Заливка
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
        cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)  # Размер шрифта
        cell.text_frame.paragraphs[0].runs[0].font.bold = False  # Жирный
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

table.cell(10, 0).merge(table.cell(10, 1))
table.cell(10, 0).text = 'ИТОГО:'
table.cell(10, 0).fill.background()  # Заливка
table.cell(10, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
table.cell(10, 0).text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
table.cell(10, 0).text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Размер шрифта
table.cell(10, 0).text_frame.paragraphs[0].runs[0].font.bold = True  # Жирный
table.cell(10, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.RIGHT  # Центр
table.cell(10, 0).vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

SumAll = 0
SumKZ = 0
SumMax = 0
for i in range(6):
    SumAll = int(table.cell(4 + i, 2).text) + SumAll
    SumKZ = int(table.cell(4 + i, 3).text) + SumKZ
    SumMax = int(table.cell(4 + i, 5).text) + SumMax

SumDZ = SumAll - SumKZ
SumItog = round(SumDZ / SumMax, 1)
if SumItog <= 1:
    SumItog = "Норма"

txt = [f'{SumAll}', f'{SumKZ}', f'{SumDZ}', f'{SumMax}', f'{SumItog}', ]
for y in range(5):
    cell = table.cell(10, y + 2)
    cell.text = txt[y]
    cell.fill.background()  # Заливка
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 32, 96)  # Цвет
    cell.text_frame.paragraphs[0].runs[0].font.name = 'Sitka Small'  # Шрифт
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Размер шрифта
    cell.text_frame.paragraphs[0].runs[0].font.bold = True  # Жирный
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER.CENTER  # Центр
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Центр

print(13)
prs.save(Name)
